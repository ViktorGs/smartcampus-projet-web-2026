"""
Génère le PowerPoint de soutenance SmartCampus.
Couvre tous les critères de la grille d'évaluation PowerPoint (25 pts).
Exécution : python build_pptx.py
Sortie    : ../SmartCampus-Presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----- Palette ---------------------------------------------------------------
NAVY      = RGBColor(0x0B, 0x2A, 0x4A)   # bandeau / titres
ACCENT    = RGBColor(0x2E, 0x86, 0xC1)   # bleu accent
TEAL      = RGBColor(0x16, 0xA0, 0x85)   # vert accent
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFB)   # fond doux
DARK_TEXT = RGBColor(0x1B, 0x26, 0x31)
GREY      = RGBColor(0x6B, 0x7A, 0x89)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
SHOTS  = ROOT / "screenshots"
OUT    = ROOT.parent / "SmartCampus-Presentation.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ----- Helpers ---------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=DARK_TEXT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = "•  " + it
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_image_fit(slide, path, x, y, max_w, max_h):
    """Ajoute une image en la rétrécissant pour entrer dans le cadre, et la centre."""
    from PIL import Image
    if not path.exists():
        return None
    im = Image.open(path)
    iw, ih = im.size
    rw = max_w / iw
    rh = max_h / ih
    r = min(rw, rh)
    w = int(iw * r)
    h = int(ih * r)
    cx = x + (max_w - w) // 2
    cy = y + (max_h - h) // 2
    return slide.shapes.add_picture(str(path), cx, cy, width=w, height=h)


def slide_header(slide, title, subtitle=None, page_num=None, total=None):
    """Bandeau supérieur sur toutes les slides internes."""
    add_rect(slide, 0, 0, SW, Inches(0.85), NAVY)
    add_rect(slide, 0, Inches(0.85), SW, Inches(0.06), ACCENT)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(11), Inches(0.6),
             title, size=24, bold=True, color=WHITE)
    if page_num is not None and total is not None:
        add_text(slide, Inches(12.0), Inches(0.25), Inches(1.0), Inches(0.4),
                 f"{page_num}/{total}", size=12, color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=GREY)


def footer(slide):
    add_rect(slide, 0, Inches(7.2), SW, Inches(0.3), LIGHT_BG)
    add_text(slide, Inches(0.4), Inches(7.22), Inches(8), Inches(0.3),
             "SmartCampus  ·  Projet Web dynamique 2026  ·  ING2 — ECE",
             size=10, color=GREY)
    add_text(slide, Inches(9.5), Inches(7.22), Inches(3.5), Inches(0.3),
             "GOUSSOT · CARMINATI · PEZE · OSORIO",
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


# ============================================================================
# Slide 1 — Page de garde
# ============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(2.9), SW, Inches(2.2), ACCENT)

# Logo ECE en haut à droite si dispo
logo = ASSETS / "ece-logo.jpg"
if logo.exists():
    add_image_fit(s, logo, Inches(11.5), Inches(0.3), Inches(1.6), Inches(1.0))

add_text(s, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
         "ING2 · Projet Web dynamique 2026", size=18, color=WHITE)
add_text(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(1.2),
         "SmartCampus", size=66, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8),
         "La gestion académique de notre époque", size=28, color=WHITE)

add_text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
         "Soutenance du projet final  ·  31 mai 2026",
         size=20, color=WHITE)
add_text(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.5),
         "Équipe : Viktor GOUSSOT  ·  Nicolas CARMINATI  ·  Louis PEZE  ·  Alexis OSORIO",
         size=16, color=WHITE)
add_text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
         "École d'ingénieurs ciblée — Filières L1→M2, semestres, ECTS, groupes TD",
         size=14, color=WHITE)


# ============================================================================
# Slide 2 — Sommaire
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "Sommaire")
add_bullets(s, Inches(0.8), Inches(1.4), Inches(6), Inches(5.5), [
    "1. Présentation générale du projet",
    "2. Équipe et répartition du travail",
    "3. Architecture finale du système",
    "4. Modèle entité-association",
    "5. Modèle relationnel",
    "6. Wireframes & storyboard",
    "7. Workflows transactionnels & UX",
], size=18)
add_bullets(s, Inches(7.0), Inches(1.4), Inches(6), Inches(5.5), [
    "8. Spécifications fonctionnelles",
    "9. Règles métier académiques",
    "10. Journal d'assistance par IA",
    "11. Rapport de compromis techniques",
    "12. Versioning Git & GitHub",
    "13. Bilans individuel et collectif",
    "14. Références et conclusion",
], size=18)
footer(s)


# ============================================================================
# Slide 3 — Présentation générale
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "1. Présentation générale")
add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
         "Plateforme web de gestion académique pour une école d'ingénieurs",
         size=18, bold=True, color=ACCENT)

# 3 colonnes
add_rect(s, Inches(0.5), Inches(1.9),  Inches(4.0), Inches(2.5), LIGHT_BG)
add_rect(s, Inches(4.7), Inches(1.9),  Inches(4.0), Inches(2.5), LIGHT_BG)
add_rect(s, Inches(8.9), Inches(1.9),  Inches(4.0), Inches(2.5), LIGHT_BG)

add_text(s, Inches(0.7), Inches(2.05), Inches(3.7), Inches(0.4),
         "Contexte ciblé", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(3.7), Inches(2.0), [
    "École d'ingénieurs",
    "Filières L1 → M2",
    "Semestres & ECTS",
    "Groupes TD",
], size=13)

add_text(s, Inches(4.9), Inches(2.05), Inches(3.7), Inches(0.4),
         "Acteurs (3 rôles)", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(4.9), Inches(2.5), Inches(3.7), Inches(2.0), [
    "Administrateur",
    "Enseignant",
    "Étudiant",
    "Accès cloisonnés (RBAC)",
], size=13)

add_text(s, Inches(9.1), Inches(2.05), Inches(3.7), Inches(0.4),
         "Périmètre fonctionnel", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(9.1), Inches(2.5), Inches(3.7), Inches(2.0), [
    "CRUD utilisateurs / cours",
    "Inscriptions, notes",
    "Emploi du temps",
    "PDF, messagerie, stats",
], size=13)

add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
         "Objectifs techniques", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(2.0), [
    "Architecture client–serveur explicite : API REST PHP 8 + SPA HTML/CSS/JS",
    "Base de données MySQL 8 structurée (PDO, requêtes préparées)",
    "Pas de framework lourd : mini-noyau MVC maison, code maîtrisé et explicable",
    "6 règles métier implémentées et testées (double inscription, capacité, prérequis, conflits, verrouillage notes)",
], size=14)
footer(s)


# ============================================================================
# Slide 4 — Équipe & répartition du travail
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "2. Équipe & répartition du travail")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "4 membres — rôles complémentaires front / back, revue croisée",
         size=14, color=GREY)

# Tableau équipe
rows, cols = 5, 3
left, top = Inches(0.5), Inches(1.75)
width, height = Inches(12.3), Inches(2.6)
tbl = s.shapes.add_table(rows, cols, left, top, width, height).table
tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(3.5)
tbl.columns[2].width = Inches(6.2)
headers = ["Membre", "Rôle principal", "Modules portés"]
data = [
    ("Viktor GOUSSOT",   "Chef de projet · Backend",
     "Noyau API, Auth & sessions, CSRF/RBAC, CRUD utilisateurs/étudiants/profs/cours"),
    ("Nicolas CARMINATI", "Backend · Données · Règles métier",
     "MCD/MR, schéma SQL, règles académiques (Academic.php), API notes/inscriptions, PDF"),
    ("Louis PEZE",        "Frontend · Pages de gestion",
     "Coquille SPA, design Bootstrap, pages étudiants/profs/cours, filtres, messagerie"),
    ("Alexis OSORIO",     "Frontend · Dashboards & EDT · Tests",
     "Tableaux de bord 3 rôles, emploi du temps, statistiques Chart.js, scénarios de test"),
]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    cell.text_frame.text = ""
    p = cell.text_frame.paragraphs[0]
    r = p.add_run(); r.text = h
    r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(14)
for i, row in enumerate(data, start=1):
    for j, v in enumerate(row):
        cell = tbl.cell(i, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BG if i % 2 else WHITE
        cell.text_frame.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = v
        r.font.size = Pt(12); r.font.color.rgb = DARK_TEXT
        if j == 0:
            r.font.bold = True

add_text(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4),
         "Organisation", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(2.0), [
    "Branches courtes par fonctionnalité, merge sur main après relecture croisée",
    "Conventions de commit (feat / fix / docs / chore) — historique lisible sur GitHub",
    "Chaque membre relit le code des autres — objectif : pouvoir expliquer tout le projet en soutenance",
    "Points d'avancement réguliers — répartition vérifiable dans les contributions Git",
], size=14)
footer(s)


# ============================================================================
# Slide 5 — Architecture finale
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "3. Architecture finale du système")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "Séparation stricte frontend / backend autour d'une API REST",
         size=14, color=GREY)

archi_png = ROOT / "Schema-Architecture.png"
if archi_png.exists():
    add_image_fit(s, archi_png, Inches(0.5), Inches(1.7), Inches(8.0), Inches(5.2))
else:
    # fallback : on dessine 3 boites
    add_rect(s, Inches(0.6), Inches(2.0), Inches(2.4), Inches(4.5), ACCENT)
    add_text(s, Inches(0.7), Inches(2.1), Inches(2.2), Inches(4.3),
             "Frontend SPA\n\n• HTML5 / CSS3\n• ES Modules\n• Bootstrap 5\n• Chart.js",
             size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.4), Inches(2.0), Inches(2.4), Inches(4.5), NAVY)
    add_text(s, Inches(3.5), Inches(2.1), Inches(2.2), Inches(4.3),
             "Backend PHP\n\n• API REST\n• Mini-MVC\n• PDO\n• Auth + CSRF\n• FPDF",
             size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.2), Inches(2.0), Inches(2.4), Inches(4.5), TEAL)
    add_text(s, Inches(6.3), Inches(2.1), Inches(2.2), Inches(4.3),
             "MySQL 8\n\n• 13 tables\n• FK contraintes\n• Index pertinents",
             size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# Légende à droite
add_text(s, Inches(8.9), Inches(1.8), Inches(4.0), Inches(0.4),
         "Choix structurants", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(8.9), Inches(2.25), Inches(4.0), Inches(5.0), [
    "Front et back même origine (Apache)",
    "Communication JSON / fetch + cookie session",
    "Pas de framework JS : ES Modules natifs",
    "Pas de framework PHP : routeur maison (~6 classes core)",
    "Vendors locaux : démo hors-ligne",
    "Validation côté client ET côté serveur",
    "RBAC central avant chaque action sensible",
], size=12)
footer(s)


# ============================================================================
# Slide 6 — MCD
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "4. Modèle entité-association (MCD)")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "13 entités, héritage users → profils, cardinalités cohérentes",
         size=14, color=GREY)

mcd_png = ROOT / "Diagramme-ER.png"
if mcd_png.exists():
    add_image_fit(s, mcd_png, Inches(0.4), Inches(1.7), Inches(9.0), Inches(5.3))

add_text(s, Inches(9.7), Inches(1.8), Inches(3.3), Inches(0.4),
         "Choix de modélisation", size=14, bold=True, color=NAVY)
add_bullets(s, Inches(9.7), Inches(2.2), Inches(3.3), Inches(5.0), [
    "users + student_profiles + teacher_profiles (héritage 1-1)",
    "Authentification unifiée, emails uniques tous rôles",
    "courses ⟷ enrollments ⟷ students (table de jointure)",
    "grades : (enrollment, assessment, value, locked)",
    "schedule_slots : créneaux horaires liés à un cours + salle + prof",
    "messages : référencent users (anti-duplication)",
], size=11)
footer(s)


# ============================================================================
# Slide 7 — Modèle relationnel
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "5. Modèle relationnel final")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "Tables principales (clés primaires soulignées · #clé étrangère)",
         size=14, color=GREY)

# Présenté en deux colonnes
left_lines = [
    ("users",            "(id, email, password_hash, role, full_name, created_at)"),
    ("student_profiles", "(#user_id, student_number, level, programme_id, group_label, ...)"),
    ("teacher_profiles", "(#user_id, office, title, biography, ...)"),
    ("programmes",       "(id, code, name, level)"),
    ("courses",          "(id, code, name, ects, semester, capacity, #teacher_id, #programme_id)"),
    ("prerequisites",    "(#course_id, #prerequisite_course_id)"),
    ("rooms",            "(id, label, capacity, building)"),
]
right_lines = [
    ("enrollments",      "(id, #student_id, #course_id, enrolled_at, status)"),
    ("assessments",      "(id, #course_id, label, type, weight, due_date)"),
    ("grades",           "(id, #enrollment_id, #assessment_id, value, locked, locked_at)"),
    ("schedule_slots",   "(id, #course_id, #room_id, weekday, start_time, end_time)"),
    ("messages",         "(id, #sender_id, #recipient_id, subject, body, sent_at, read_at)"),
    ("notifications",    "(id, #user_id, type, payload, created_at, read_at)"),
    ("audit_logs",       "(id, #user_id, action, target, created_at)"),
]
def render_tables(slide, x, items):
    y = Inches(1.75)
    for name, schema in items:
        add_rect(slide, x, y, Inches(6.2), Inches(0.7), LIGHT_BG)
        add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.0), Inches(0.3),
                 name, size=13, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.15), y + Inches(0.32), Inches(5.9), Inches(0.35),
                 schema, size=11, color=DARK_TEXT)
        y += Inches(0.75)

render_tables(s, Inches(0.4),  left_lines)
render_tables(s, Inches(6.75), right_lines)
footer(s)


# ============================================================================
# Slide 8 — Wireframes & storyboard
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "6. Wireframes & storyboard")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "Maquettes principales et parcours utilisateurs validés en amont",
         size=14, color=GREY)

shots = [
    (SHOTS / "login.png",              "Connexion"),
    (SHOTS / "student-dashboard.png",  "Dashboard étudiant"),
    (SHOTS / "teacher-dashboard.png",  "Dashboard enseignant"),
    (SHOTS / "admin-dashboard.png",    "Dashboard admin"),
    (SHOTS / "student-catalogue.png",  "Catalogue cours (filtres/tri)"),
    (SHOTS / "admin-schedule.png",     "Emploi du temps"),
]
# grille 3x2
gx, gy = Inches(0.4), Inches(1.7)
cw, ch = Inches(4.2), Inches(2.5)
for i, (p, caption) in enumerate(shots):
    col = i % 3
    row = i // 3
    x = gx + col * (cw + Inches(0.1))
    y = gy + row * (ch + Inches(0.3))
    add_rect(s, x, y, cw, ch, LIGHT_BG)
    add_image_fit(s, p, x + Inches(0.08), y + Inches(0.08),
                  cw - Inches(0.16), ch - Inches(0.5))
    add_text(s, x, y + ch - Inches(0.35), cw, Inches(0.3),
             caption, size=11, color=NAVY, align=PP_ALIGN.CENTER, bold=True)
footer(s)


# ============================================================================
# Slide 9 — Workflows transactionnels
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "7. Workflows transactionnels & UX")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "Parcours critiques : inscription · saisie de note · détection conflit",
         size=14, color=GREY)

def workflow(slide, y, title, steps, color):
    add_text(slide, Inches(0.5), y, Inches(12.3), Inches(0.35),
             title, size=15, bold=True, color=color)
    sx = Inches(0.5)
    sw = Inches(2.3)
    sh = Inches(0.7)
    sy = y + Inches(0.4)
    for i, step in enumerate(steps):
        x = sx + i * (sw + Inches(0.2))
        add_rect(slide, x, sy, sw, sh, color)
        add_text(slide, x + Inches(0.1), sy + Inches(0.1), sw - Inches(0.2), sh - Inches(0.2),
                 step, size=11, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

workflow(s, Inches(1.6),
         "Inscription d'un étudiant à un cours",
         ["Étudiant\nclique S'inscrire",
          "API: vérifie\nrôle + CSRF",
          "Règles métier:\ndoublon, capacité,\nprérequis, conflit",
          "INSERT enrollments\n(transaction)",
          "UI: toast OK\n+ refresh listes"],
         ACCENT)

workflow(s, Inches(3.7),
         "Saisie + verrouillage d'une note",
         ["Enseignant\nsaisit note",
          "Validation\n(0 ≤ note ≤ 20)",
          "INSERT/UPDATE\ngrades",
          "Validation finale\n→ locked=1",
          "Modif refusée\n(blocage métier)"],
         NAVY)

workflow(s, Inches(5.8),
         "Détection conflit emploi du temps",
         ["Création\ncréneau",
          "SELECT slots\nWHERE start<? AND end>?",
          "Conflit salle\nou enseignant ?",
          "Erreur 409\n+ message clair",
          "Sinon INSERT\nschedule_slots"],
         TEAL)
footer(s)


# ============================================================================
# Slide 10 — Spécifications fonctionnelles
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "8. Spécifications fonctionnelles")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "Fonctionnalités principales par rôle utilisateur",
         size=14, color=GREY)

rows, cols = 4, 3
tbl = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.75),
                         Inches(12.3), Inches(2.5)).table
for c, w in zip(range(3), [Inches(2.5), Inches(4.9), Inches(4.9)]):
    tbl.columns[c].width = w
hdr = ["Rôle", "Espace dédié", "Actions clés"]
for j, h in enumerate(hdr):
    cell = tbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    cell.text_frame.text = ""
    r = cell.text_frame.paragraphs[0].add_run()
    r.text = h; r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(13)

data = [
    ("Étudiant",      "Tableau de bord, catalogue, mes cours, notes, EDT, messagerie",
     "S'inscrire / se désinscrire · consulter notes & EDT · télécharger PDF du relevé"),
    ("Enseignant",    "Dashboard, mes cours, saisie notes, EDT, messagerie",
     "Saisir/modifier notes · valider (verrouiller) · gérer cours · écrire aux étudiants"),
    ("Administrateur","Gestion users, étudiants, profs, cours, EDT, stats globales",
     "CRUD complet · création de cours · planning salles · statistiques"),
]
for i, row in enumerate(data, 1):
    for j, v in enumerate(row):
        cell = tbl.cell(i, j); cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BG if i % 2 else WHITE
        cell.text_frame.text = ""
        r = cell.text_frame.paragraphs[0].add_run()
        r.text = v; r.font.size = Pt(11); r.font.color.rgb = DARK_TEXT
        if j == 0:
            r.font.bold = True

add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
         "Justification des choix fonctionnels", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.2), [
    "Une SPA unique : pas de rechargement, expérience fluide, code organisé en modules ES",
    "API REST stateless côté logique métier, session côté authentification (même origine)",
    "Recherche + filtres + tri pour les cours : poste à 4 pts dans la grille — UX prioritaire",
    "PDF / messagerie / statistiques implémentés en bonus pour démontrer la maîtrise",
], size=13)
footer(s)


# ============================================================================
# Slide 11 — Règles métier
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "9. Règles métier académiques")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "6 règles métier implémentées et centralisées dans core/Academic.php",
         size=14, color=GREY)

rules = [
    ("Double inscription",
     "Un étudiant ne peut s'inscrire qu'une fois à un même cours actif.",
     "UNIQUE(student_id, course_id) + pré-check côté contrôleur"),
    ("Capacité maximale",
     "Refus si le cours a atteint sa capacité (champ courses.capacity).",
     "COUNT(enrollments) avant INSERT, dans la même transaction"),
    ("Prérequis académiques",
     "Inscription bloquée si un prérequis du cours n'est pas validé.",
     "JOIN prerequisites + grades.locked AND value ≥ seuil"),
    ("Conflit horaire étudiant",
     "Refus si le créneau chevauche un cours déjà suivi par l'étudiant.",
     "SQL : a.start < b.end AND a.end > b.start (sur slots × enrollments)"),
    ("Conflit ressource (salle / prof)",
     "Refus de planifier un slot si la salle ou l'enseignant est déjà pris.",
     "Même requête de chevauchement, appliquée sur la ressource"),
    ("Verrouillage des notes",
     "Une note validée (locked=1) ne peut plus être modifiée.",
     "Garde côté contrôleur + UI désactivée + audit_logs"),
]

for i, (name, desc, impl) in enumerate(rules):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    y = Inches(1.75) + row * Inches(1.7)
    add_rect(s, x, y, Inches(6.0), Inches(1.55), LIGHT_BG)
    add_rect(s, x, y, Inches(0.15), Inches(1.55), TEAL)
    add_text(s, x + Inches(0.3), y + Inches(0.08), Inches(5.5), Inches(0.35),
             f"{i+1}. {name}", size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(0.45), Inches(5.5), Inches(0.5),
             desc, size=11, color=DARK_TEXT)
    add_text(s, x + Inches(0.3), y + Inches(1.02), Inches(5.5), Inches(0.45),
             "Impl. : " + impl, size=10, color=GREY)
footer(s)


# ============================================================================
# Slide 12 — Journal IA (1/2)
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "10. Journal d'assistance par IA (1/2)")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "Outils utilisés et tâches assistées — l'IA comme assistante, pas comme substitut",
         size=14, color=GREY)

add_text(s, Inches(0.5), Inches(1.85), Inches(6), Inches(0.4),
         "Outils utilisés", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(2.3), Inches(6), Inches(4.5), [
    "Assistant IA (Claude / ChatGPT) — conception, génération de squelettes, explications, débogage",
    "GitHub Copilot — auto-complétion ponctuelle dans VS Code (composants UI répétitifs)",
    "Documentation officielle — PHP (PDO, password_hash), MDN (fetch, ES Modules), Bootstrap, FPDF",
    "StackOverflow — vérifications croisées sur les questions SQL non triviales",
], size=12)

add_text(s, Inches(6.8), Inches(1.85), Inches(6), Inches(0.4),
         "Tâches assistées par IA", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(6.8), Inches(2.3), Inches(6), Inches(4.5), [
    "Modèle entité-association — discussion des cardinalités (héritage users→profils)",
    "Noyau backend — routeur REST, réponses JSON, validation",
    "Règles métier — requête SQL de chevauchement horaire (start<? AND end>?)",
    "Composants frontend répétitifs — modales, tableaux dynamiques",
    "Débogage ponctuel — génération PDF, encodage CSRF",
    "Rédaction documentaire — squelettes des fichiers docs/ ensuite relus et adaptés",
], size=12)
footer(s)


# ============================================================================
# Slide 13 — Journal IA (2/2) — analyse critique + limites
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "10. Journal IA (2/2) — analyse critique & limites")

add_text(s, Inches(0.5), Inches(1.35), Inches(6), Inches(0.4),
         "Réponses utiles", size=16, bold=True, color=TEAL)
add_bullets(s, Inches(0.5), Inches(1.8), Inches(6), Inches(2.4), [
    "Structure mini-MVC sans framework : claire et facile à expliquer",
    "Requête SQL de chevauchement horaire : correcte et concise",
    "Rappels sécurité : bcrypt, requêtes préparées, HttpOnly, CSRF",
    "Suggestions d'organisation des dossiers (core/, controllers/)",
], size=12)

add_text(s, Inches(6.8), Inches(1.35), Inches(6), Inches(0.4),
         "Réponses corrigées / écartées", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(6.8), Inches(1.8), Inches(6), Inches(2.4), [
    "FPDF sans dossier font/ → PDF corrompu ; ajout du dossier + désactivation des warnings",
    "Proposition Laravel : écartée (sujet déconseille les générateurs lourds)",
    "Auth JWT proposée : remplacée par sessions PHP (même origine, plus simple)",
    "Validation uniquement côté client : doublée côté serveur (client contournable)",
    "Identifiants DB en dur : déplacés dans config.local.php non versionné",
], size=12)

add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
         "Limites observées (recul critique)", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.2), [
    "Code plausible mais incomplet : dépendances oubliées (polices FPDF) → toujours tester réellement",
    "Manque de contexte : ne connaît ni notre version PHP exacte ni notre mot de passe MySQL → adaptation manuelle",
    "Tendance au sur-dimensionnement : propose souvent Docker, Redis, Vite… → écarté pour rester pédagogique",
    "Notre démarche : lire + comprendre + adapter + tester + savoir expliquer (chacun doit pouvoir défendre le code)",
], size=12)
footer(s)


# ============================================================================
# Slide 14 — Rapport de compromis techniques
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "11. Rapport de compromis techniques")

add_text(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
         "Difficultés rencontrées", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(1.75), Inches(6), Inches(2.5), [
    "FPDF cassé par un warning PHP avant le binaire (fix : ajout font/ + display_errors=0)",
    "Détection des conflits horaires en SQL : exprimer le chevauchement et distinguer conflit étudiant vs ressource",
    "Routage sans framework : segments dynamiques {id} compatibles Apache .htaccess et serveur PHP intégré",
    "Cohérence cache local frontend / base après mutation : refresh ciblé après chaque action",
], size=12)

add_text(s, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4),
         "Compromis assumés", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(6.8), Inches(1.75), Inches(6), Inches(2.5), [
    "Sessions PHP plutôt que JWT : moins flexible multi-clients, plus simple en même origine",
    "Mini-MVC maison plutôt qu'un framework : réinvention de briques basiques, en échange compréhension totale",
    "ES Modules sans bundler : pas d'optimisation, mais zéro étape de build",
    "Vendors locaux (Bootstrap, Chart.js) : ~900 Ko versionnés, démo hors-ligne garantie",
], size=12)

add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
         "Limites actuelles & évolutions possibles", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.2), [
    "Pas de pagination sur les listes : suffisant en démo, à ajouter pour des effectifs réels",
    "Sécurité « pédagogique » : pas de rate limiting, pas de HTTPS local, politique de mot de passe simple",
    "Emploi du temps hebdo générique : pas de gestion calendaire (dates, semaines A/B, vacances)",
    "Évolutions : pagination serveur, présences/QR code, export Excel, calendrier interactif, tests automatisés",
], size=12)
footer(s)


# ============================================================================
# Slide 15 — Versioning Git & GitHub
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "12. Versioning Git & GitHub")
add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4),
         "Dépôt : github.com/ViktorGs/smartcampus-projet-web-2026",
         size=14, color=GREY)

add_text(s, Inches(0.5), Inches(1.75), Inches(6), Inches(0.4),
         "Historique de progression", size=16, bold=True, color=ACCENT)
commits = [
    ("25 mai", "chore: initialisation du projet (README, gitignore)"),
    ("25 mai", "feat(db): schéma relationnel + données de démonstration"),
    ("26 mai", "feat(api): noyau backend (PDO, routeur REST, auth, CSRF, validation)"),
    ("27 mai", "feat(api): CRUD utilisateurs, étudiants, enseignants, cours"),
    ("28 mai", "feat(api): inscriptions, notes, emploi du temps + règles métier"),
    ("29 mai", "feat(api): bonus — relevé PDF, messagerie, notifications, dashboards, stats"),
    ("30 mai", "feat(front): coquille SPA, design Bootstrap, noyau JS"),
    ("30 mai", "feat(front): pages (dashboards, étudiants, cours, notes, EDT, messagerie, stats)"),
    ("31 mai", "docs: spécifications, MCD, architecture, wireframes, rapports + livrable final"),
]
y = Inches(2.25)
for d, msg in commits:
    add_rect(s, Inches(0.5), y, Inches(0.9), Inches(0.32), ACCENT)
    add_text(s, Inches(0.5), y, Inches(0.9), Inches(0.32),
             d, size=10, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), y + Inches(0.04), Inches(11.0), Inches(0.3),
             msg, size=11, color=DARK_TEXT)
    y += Inches(0.42)

add_text(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(0.4),
         "Organisation GitHub", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), y + Inches(0.55), Inches(12.3), Inches(1.5), [
    "Conventions de commit : feat / fix / docs / chore — historique lisible",
    "Branche main maintenue stable ; branches courtes par fonctionnalité",
    "Pas de commit massif de fin : progression jour par jour visible dans `git log`",
], size=12)
footer(s)


# ============================================================================
# Slide 16 — Bilan individuel
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "13. Bilan individuel")

people = [
    ("Viktor GOUSSOT", ACCENT, [
        "Réalisé : noyau backend, auth, CSRF, RBAC, CRUD principal",
        "Appris : architecture MVC sans framework, sécurité applicative",
        "Difficulté : routage segments dynamiques compatible Apache + PHP intégré",
    ]),
    ("Nicolas CARMINATI", NAVY, [
        "Réalisé : MCD/MR, schéma SQL, règles métier centralisées, PDF",
        "Appris : modélisation académique, requêtes SQL non triviales",
        "Difficulté : génération PDF cassée par un warning PHP",
    ]),
    ("Louis PEZE", TEAL, [
        "Réalisé : coquille SPA, pages de gestion, filtres, messagerie",
        "Appris : ES Modules natifs, design system Bootstrap",
        "Difficulté : cohérence cache local après mutation simultanée",
    ]),
    ("Alexis OSORIO", ACCENT, [
        "Réalisé : dashboards 3 rôles, EDT, statistiques Chart.js, tests",
        "Appris : visualisation de données, UI contextuelle par rôle",
        "Difficulté : lisibilité de l'emploi du temps sur mobile",
    ]),
]
for i, (name, color, bullets) in enumerate(people):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    y = Inches(1.35) + row * Inches(2.85)
    add_rect(s, x, y, Inches(6.0), Inches(2.7), LIGHT_BG)
    add_rect(s, x, y, Inches(0.15), Inches(2.7), color)
    add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(5.5), Inches(0.4),
             name, size=15, bold=True, color=NAVY)
    add_bullets(s, x + Inches(0.3), y + Inches(0.6), Inches(5.5), Inches(2.0),
                bullets, size=11)
footer(s)


# ============================================================================
# Slide 17 — Bilan collectif
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "14. Bilan collectif")

add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
         "Ce qui a bien fonctionné", size=18, bold=True, color=TEAL)
add_bullets(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.8), [
    "Séparation claire frontend / backend — chacun a pu avancer en parallèle",
    "Mini-noyau MVC explicite : chaque membre peut le lire et l'expliquer en soutenance",
    "Choix d'un contexte précis (école d'ingénieurs) qui a guidé toutes les décisions",
    "Conventions de commit partagées : historique Git lisible et exploitable",
], size=13)

add_text(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.4),
         "Ce qui a été difficile", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(1.4), [
    "Règles métier transverses (conflits horaires étudiant vs ressource)",
    "Gestion fine des permissions : 3 rôles × N actions",
    "Renoncer aux fonctionnalités séduisantes hors-périmètre (QR code, WebSocket temps réel)",
], size=13)

add_text(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.4),
         "Ce que l'on referait différemment", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.1), [
    "Mettre en place des tests automatisés dès le début (PHPUnit + Playwright)",
    "Commencer par la maquette des dashboards (point de convergence)",
    "Prévoir la pagination dès le schéma initial",
], size=13)
footer(s)


# ============================================================================
# Slide 18 — Références
# ============================================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "15. Références et ressources")

add_text(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4),
         "Bibliothèques & frameworks", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(1.85), Inches(6), Inches(3.5), [
    "Bootstrap 5.3 — design system",
    "Bootstrap Icons — pictogrammes",
    "Chart.js 4 — graphiques (statistiques)",
    "FPDF 1.86 — génération PDF du relevé",
    "PHP 8.2 — backend",
    "MySQL 8.4 — base de données",
], size=12)

add_text(s, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4),
         "Documentation & outils", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(6.8), Inches(1.85), Inches(6), Inches(3.5), [
    "php.net — PDO, password_hash, sessions",
    "MDN Web Docs — fetch, ES Modules, FormData",
    "getbootstrap.com — composants UI",
    "WAMP Server — environnement local",
    "VS Code + Live Server — édition / preview",
    "Git + GitHub — versioning et collaboration",
], size=12)

add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
         "Outils d'IA générative", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.4), [
    "Claude (Anthropic) — conception, génération de squelettes, débogage, rédaction documentaire",
    "ChatGPT (OpenAI) — vérifications croisées, explications de mécanismes",
    "GitHub Copilot — auto-complétion ponctuelle dans VS Code",
], size=12)
footer(s)


# ============================================================================
# Slide 19 — Démonstration / merci
# ============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.0), SW, Inches(1.5), ACCENT)
add_text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.6),
         "SmartCampus", size=48, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.5),
         "Place à la démonstration",
         size=28, color=WHITE)

add_text(s, Inches(0.8), Inches(3.15), Inches(11.7), Inches(0.6),
         "Merci de votre attention.",
         size=28, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.5),
         "Questions ?", size=22, color=WHITE)

add_text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.4),
         "github.com/ViktorGs/smartcampus-projet-web-2026",
         size=16, color=WHITE)
add_text(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
         "Équipe : Viktor GOUSSOT · Nicolas CARMINATI · Louis PEZE · Alexis OSORIO",
         size=14, color=WHITE)


# ----- Sauvegarde ------------------------------------------------------------
prs.save(OUT)
print(f"OK : {OUT}  ({OUT.stat().st_size // 1024} Ko)")
